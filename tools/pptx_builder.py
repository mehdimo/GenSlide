"""
PPTX builder node for the GenSlide LangGraph pipeline.

Responsibilities:
    - Take state["slides"] (List[SlideContent]) produced by the content agent
    - Assemble a polished .pptx file using python-pptx
    - Write the output path into state["pptx_path"]
    - Write a human-readable error into state["error"] on failure

Slide layout:
    Slide 0  — Title slide  (deck title + subtitle)
    Slides 1…N — Content slides  (title bar + bullet body + speaker notes)

Design:
    Theme       : "Midnight Executive" — navy (#1E2761) dominant,
                  ice-blue (#CADCFC) accent, white text
    Structure   : Dark title/outro sandwich, light content slides
    Motif       : Thick left-side accent bar on every content slide
                  (carries the navy color across all slides consistently)
    Typography  : Calibri throughout; large title, readable body (20pt)
"""

import logging
import os
from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from graph.state import AgentState, SlideContent

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Design constants  (Midnight Executive palette)
# ---------------------------------------------------------------------------

# Colors (no # prefix for RGBColor — pass as 0xRRGGBB int or use from_string)
C_NAVY      = RGBColor(0x1E, 0x27, 0x61)   # dominant background
C_ICE_BLUE  = RGBColor(0xCA, 0xDC, 0xFC)   # accent / highlights
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)   # near-black for body on light bg
C_LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xFF)   # very pale blue-white for content slides
C_ACCENT_BAR= RGBColor(0x1E, 0x27, 0x61)   # left-side motif bar

# Slide dimensions (16:9 — 10" × 5.625")
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Layout grid (inches)
MARGIN      = Inches(0.4)
BAR_W       = Inches(0.12)       # left accent bar width
CONTENT_X   = Inches(0.65)       # text starts after bar + gap
CONTENT_W   = Inches(9.1)        # width of text regions
TITLE_Y     = Inches(0.28)
TITLE_H     = Inches(0.72)
BODY_Y      = Inches(1.15)
BODY_H      = Inches(3.95)
FOOTER_Y    = Inches(5.15)
FOOTER_H    = Inches(0.35)

OUTPUT_DIR  = Path("frontend/generated")


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _rgb(color: RGBColor):
    """Return the six-char hex string (no #) used in XML attributes."""
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _set_bg(slide, color: RGBColor):
    """Fill the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color: RGBColor, transparency: int = 0):
    """Add a filled rectangle shape with no border."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.fill.background()          # no border
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    if transparency:
        # transparency 0–100 → stored as 0–100000 in pptx XML
        fill.fore_color._element.attrib["lumMod"] = str(100000 - transparency * 1000)
    return shape


def _add_text_box(
    slide,
    text: str,
    x, y, w, h,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_face: str = "Calibri",
    word_wrap: bool = True,
):
    """Add a single-paragraph text box."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_face
    return txBox


def _add_bullet_body(slide, bullets: list[str], x, y, w, h):
    """
    Add a text frame with proper PptxGenJS-style bullets.
    Each bullet is its own paragraph to ensure correct rendering.
    """
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = 0
        # Enable bullet via XML (python-pptx doesn't expose bullets natively)
        pPr = p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "▸")

        # Bullet font size and color
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", "Arial")

        # Space before each bullet
        spcBef = etree.SubElement(pPr, qn("a:spcBef"))
        spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
        spcPts.set("val", "160")   # 8pt spacing above each bullet

        run = p.add_run()
        run.text = bullet_text
        run.font.size = Pt(20)
        run.font.color.rgb = C_DARK_TEXT
        run.font.name = "Calibri"

    return txBox


def _add_speaker_notes(slide, notes_text: str):
    """Write speaker notes to the notes placeholder."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def _add_footer(slide, slide_num: int, total: int, deck_title: str):
    """Add a subtle footer bar with slide number and deck title."""
    # Footer background bar
    _add_rect(slide, 0, FOOTER_Y, SLIDE_W, FOOTER_H, C_NAVY)

    # Deck title (left)
    _add_text_box(
        slide, deck_title,
        CONTENT_X, FOOTER_Y, Inches(6), FOOTER_H,
        font_size=9, color=C_ICE_BLUE,
        align=PP_ALIGN.LEFT,
    )
    # Slide number (right)
    _add_text_box(
        slide, f"{slide_num} / {total}",
        Inches(8.5), FOOTER_Y, Inches(1.3), FOOTER_H,
        font_size=9, color=C_ICE_BLUE,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _build_title_slide(prs: Presentation, deck_title: str, subtitle: str):
    """
    Slide 0 — full navy background, large white title, ice-blue subtitle.
    Accent: horizontal rule between title and subtitle.
    """
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, C_NAVY)

    # Decorative left bar (motif — thicker on title slide)
    _add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, C_ICE_BLUE)

    # Title
    _add_text_box(
        slide, deck_title,
        Inches(0.55), Inches(1.6), Inches(8.8), Inches(1.4),
        font_size=44, color=C_WHITE, bold=True,
        align=PP_ALIGN.LEFT,
    )

    # Horizontal rule
    rule = slide.shapes.add_shape(1, Inches(0.55), Inches(3.15), Inches(3.5), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = C_ICE_BLUE
    rule.line.fill.background()

    # Subtitle
    _add_text_box(
        slide, subtitle,
        Inches(0.55), Inches(3.35), Inches(8.8), Inches(0.8),
        font_size=18, color=C_ICE_BLUE,
        align=PP_ALIGN.LEFT,
    )

    # Date (bottom-right)
    date_str = datetime.today().strftime("%B %Y")
    _add_text_box(
        slide, date_str,
        Inches(7.5), Inches(5.1), Inches(2.3), Inches(0.4),
        font_size=10, color=C_ICE_BLUE,
        align=PP_ALIGN.RIGHT,
    )

    return slide


def _build_content_slide(
    prs: Presentation,
    slide_data: SlideContent,
    slide_num: int,
    total: int,
    deck_title: str,
):
    """
    Content slides — pale blue-white background, navy title bar at top,
    left accent bar (motif), bullet body, speaker notes, footer.
    """
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, C_LIGHT_BG)

    # Left accent bar (motif — carries across all content slides)
    _add_rect(slide, 0, 0, BAR_W, SLIDE_H, C_ACCENT_BAR)

    # Title bar background
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), C_NAVY)

    # Slide title text
    _add_text_box(
        slide, slide_data["title"],
        CONTENT_X, TITLE_Y, CONTENT_W, TITLE_H,
        font_size=28, color=C_WHITE, bold=True,
        align=PP_ALIGN.LEFT,
    )

    # Slide number badge (top-right inside title bar)
    _add_text_box(
        slide, str(slide_num),
        Inches(9.3), Inches(0.2), Inches(0.5), Inches(0.6),
        font_size=22, color=C_ICE_BLUE, bold=True,
        align=PP_ALIGN.CENTER,
    )

    # Bullet body
    bullets = slide_data.get("bullets") or []
    if bullets:
        _add_bullet_body(slide, bullets, CONTENT_X, BODY_Y, CONTENT_W, BODY_H)

    # Footer
    _add_footer(slide, slide_num, total, deck_title)

    # Speaker notes
    notes = slide_data.get("speaker_notes", "")
    if notes:
        _add_speaker_notes(slide, notes)

    return slide


def _build_outro_slide(prs: Presentation, deck_title: str):
    """
    Final slide — mirrors title slide style with a 'Thank you' message.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    _set_bg(slide, C_NAVY)

    _add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, C_ICE_BLUE)

    _add_text_box(
        slide, "Thank you",
        Inches(0.55), Inches(1.8), Inches(8.8), Inches(1.2),
        font_size=44, color=C_WHITE, bold=True,
        align=PP_ALIGN.LEFT,
    )

    rule = slide.shapes.add_shape(1, Inches(0.55), Inches(3.1), Inches(2.5), Pt(2))
    rule.fill.solid()
    rule.fill.fore_color.rgb = C_ICE_BLUE
    rule.line.fill.background()

    _add_text_box(
        slide, deck_title,
        Inches(0.55), Inches(3.3), Inches(8.8), Inches(0.7),
        font_size=18, color=C_ICE_BLUE,
        align=PP_ALIGN.LEFT,
    )

    return slide


# ---------------------------------------------------------------------------
# Main assembler
# ---------------------------------------------------------------------------

def _assemble_pptx(slides: list[SlideContent], iteration: int) -> str:
    """
    Build and save the full presentation.
    Returns the output file path.
    """
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Derive deck title from the first slide's title
    deck_title = slides[0]["title"] if slides else "Presentation"
    subtitle    = "Generated by GenSlide · AI-powered presentation builder"
    total_content = len(slides)

    # Title slide
    _build_title_slide(prs, deck_title, subtitle)

    # Content slides
    for i, slide_data in enumerate(slides, start=1):
        _build_content_slide(prs, slide_data, i, total_content, deck_title)

    # Outro slide
    _build_outro_slide(prs, deck_title)

    # Save
    filename = f"deck_v{iteration}_{datetime.today().strftime('%Y%m%d_%H%M%S')}.pptx"
    path = OUTPUT_DIR / filename
    prs.save(str(path))
    logger.info("PPTX saved → %s", path)
    return str(path)


# ---------------------------------------------------------------------------
# LangGraph node
# ---------------------------------------------------------------------------

def build_pptx_node(state: AgentState) -> AgentState:
    """
    LangGraph node: assemble the .pptx from state["slides"].

    Reads:
        state["slides"]     — List[SlideContent] from content_agent
        state["iteration"]  — used in the output filename

    Writes:
        state["pptx_path"]  — path to the saved file
        state["error"]      — error message on failure (else None)
    """
    slides = state.get("slides") or []

    if not slides:
        state["error"] = "No slides found in state. Content agent may have failed."
        state["pptx_path"] = None
        return state

    try:
        path = _assemble_pptx(slides, state.get("iteration", 1))
        state["pptx_path"] = path
        state["error"] = None
        logger.info("build_pptx_node: deck saved at %s", path)

    except Exception as exc:
        logger.exception("build_pptx_node failed: %s", exc)
        state["pptx_path"] = None
        state["error"] = f"PPTX generation failed: {exc}"

    return state