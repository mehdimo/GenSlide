import os

from pptx import Presentation
from pptx.util import Inches

class SlideDeck:

    def __init__(self, output_folder="generated"):
        self.prs = Presentation()
        self.output_folder = output_folder

    def add_slide(self, slide_data):
        prs = self.prs
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        # Title
        title_shape = shapes.title
        title_shape.text = slide_data.get("title_text", "")

        # Body
        if "text" in slide_data:
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            for bullet in slide_data.get("text", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

                if "p1" in slide_data:
                    p = tf.add_paragraph()
                    p.text = slide_data.get("p1")
                    p.level = 1

        if "img_path" in slide_data:
            cur_left = 6
            for img_path in slide_data.get("img_path", []):
                top = Inches(2)
                left = Inches(cur_left)
                height = Inches(4)
                pic = slide.shapes.add_picture(img_path, left, top, height=height)
                cur_left += 1

    def add_title_slide(self, title_page_data):
        # title slide
        prs = self.prs
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        if "title_text" in title_page_data:
            title.text = title_page_data.get("title_text")
        if "subtitle_text" in title_page_data:
            subtitle.text = title_page_data.get("subtitle_text")

    def create_presentation(self, title_slide_info, slide_pages_data=[]):
        try:
            file_name = title_slide_info.get("title_text").\
                lower().replace(",", "").replace(" ", "-")
            file_name += ".pptx"
            file_name = os.path.join(self.output_folder, file_name)
            self.add_title_slide(title_slide_info)
            for slide_data in slide_pages_data:
                self.add_slide(slide_data)

            self.prs.save(file_name)
            return file_name
        except Exception as e:
            raise e
